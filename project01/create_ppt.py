#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate Nebula Studio project presentation PPT."""

import sys
print("Starting...", flush=True)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    print("Imports OK", flush=True)
except Exception as e:
    print(f"Import error: {e}", flush=True)
    sys.exit(1)

# ===== Color Theme =====
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x4F, 0xC3, 0xF7)
ACCENT_PURPLE = RGBColor(0xBB, 0x86, 0xFC)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)
ACCENT_ORANGE = RGBColor(0xFF, 0x91, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB0, 0xC0)
DARK_CARD = RGBColor(0x16, 0x21, 0x3E)
MEDIUM_BG = RGBColor(0x1F, 0x29, 0x4E)
ACCENT_RED = RGBColor(0xFF, 0x55, 0x55)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, color=DARK_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE,
                bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_section_slide(prs, number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, MEDIUM_BG)
    add_textbox(slide, Inches(1.5), Inches(1.5), Inches(3), Inches(2.5),
                str(number).zfill(2), font_size=96, color=ACCENT_PURPLE, bold=True, font_name='Arial')
    add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(1.5),
                title, font_size=42, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(5.8), Inches(2.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()
    return slide


# ==================== SLIDES ====================

# --- Slide 1: Cover ---
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1)
# decorative circles
for x_off, sz, clr in [(2.0, 0.6, ACCENT_PURPLE), (8.5, 0.35, ACCENT_BLUE), (10.8, 0.45, ACCENT_GREEN)]:
    c = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_off), Inches(1.2), Inches(sz), Inches(sz))
    c.fill.solid()
    c.fill.fore_color.rgb = clr
    c.line.fill.background()

add_textbox(slide1, Inches(1.8), Inches(2.5), Inches(10), Inches(1.5),
            'Nebula Studio', font_size=56, bold=True, font_name='Arial')
add_textbox(slide1, Inches(1.8), Inches(3.8), Inches(10), Inches(1.0),
            '文图互转主题设计系统', font_size=32, color=ACCENT_BLUE)
line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.8), Inches(5.0), Inches(3.5), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_PURPLE
line.line.fill.background()
add_textbox(slide1, Inches(1.8), Inches(5.3), Inches(10), Inches(0.5),
            '项目终期汇报', font_size=22, color=LIGHT_GRAY)
add_textbox(slide1, Inches(1.8), Inches(5.9), Inches(10), Inches(0.5),
            '2026年5月  |  版本 1.0', font_size=16, color=LIGHT_GRAY)

print("Slide 1 done", flush=True)

# --- Slide 2: TOC ---
add_section_slide(prs, 0, '目  录')
slide_toc = prs.slides[-1]
toc_items = [
    ('01', '项目概述', '背景、目标与痛点'),
    ('02', '系统架构', '技术栈与架构设计'),
    ('03', '功能模块', '六大核心功能'),
    ('04', 'AI核心', '策略模式 & ComfyUI集成'),
    ('05', '权限体系', 'RBAC三级角色'),
    ('06', '数据库设计', '7张核心表'),
    ('07', '前端展示', 'Vue 3 界面一览'),
    ('08', '部署管理', '开发进度与部署方案'),
]
for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(1.2 + i * 0.75)
    add_textbox(slide_toc, Inches(3.0), y, Inches(0.8), Inches(0.6),
                num, font_size=28, bold=True, color=ACCENT_BLUE, font_name='Arial')
    add_textbox(slide_toc, Inches(4.0), y, Inches(3.5), Inches(0.6),
                title, font_size=24, bold=True)
    add_textbox(slide_toc, Inches(7.5), y, Inches(4.0), Inches(0.6),
                desc, font_size=16, color=LIGHT_GRAY)

print("Slide 2 done", flush=True)

# --- Slide 3: Project Overview ---
add_section_slide(prs, 1, '项目概述')
slide3 = prs.slides[-1]

cards = [
    ('🎯 项目定位', '文图互转主题设计平台，企业级AI辅助创作工具。支持文本生成图像（Text2Image）与图像理解（Image2Text），提供完整的风格管理、LoRA微调训练链路。'),
    ('💡 核心价值', '创意可复用，风格可积累。通过主题风格管理 + LoRA微调实现品牌视觉一致性。AI辅助将设计师产出效率提升10倍以上，降低重复劳动成本。'),
    ('⚡ 解决痛点', '效率瓶颈 → AI辅助生成\n一致性难题 → 风格管理与复用\n定制化成本 → LoRA模型训练\n多角色协作 → 三级权限体系'),
]
for i, (title, desc) in enumerate(cards):
    left = Inches(1.5 + i * 3.7)
    add_rounded_rect(slide3, left, Inches(1.5), Inches(3.3), Inches(5.0), DARK_CARD)
    add_textbox(slide3, left + Inches(0.25), Inches(1.7), Inches(2.8), Inches(0.6),
                title, font_size=20, bold=True, color=ACCENT_BLUE if i == 0 else ACCENT_GREEN if i == 1 else ACCENT_ORANGE)
    add_textbox(slide3, left + Inches(0.25), Inches(2.4), Inches(2.8), Inches(3.8),
                desc, font_size=14, color=LIGHT_GRAY)

print("Slide 3 done", flush=True)

# --- Slide 4: Architecture ---
add_section_slide(prs, 2, '系统架构')
slide4 = prs.slides[-1]

layers = [
    ('前端层', 'Vue 3 + Vite + Element Plus + Pinia  (端口 3000)', ACCENT_BLUE),
    ('网关层', 'Spring Security + JWT 认证 + BCrypt密码哈希', ACCENT_GREEN),
    ('业务层', '文生图 · 图生文 · 风格管理 · LoRA训练 · 用户管理 · 系统管理', ACCENT_PURPLE),
    ('AI服务层', 'Text2ImageProvider / Image2TextProvider (策略模式)  |  ComfyUI / Doubao', ACCENT_ORANGE),
    ('数据层', 'MySQL 8.0 / H2 + Redis + MinIO + Kafka  |  MyBatis-Plus ORM', ACCENT_RED),
]
for i, (name, tech, color) in enumerate(layers):
    y = Inches(1.6 + i * 1.1)
    bar = add_rounded_rect(slide4, Inches(1.5), y, Inches(10.3), Inches(0.85), DARK_CARD)
    ind = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), y, Inches(0.08), Inches(0.85))
    ind.fill.solid()
    ind.fill.fore_color.rgb = color
    ind.line.fill.background()
    add_textbox(slide4, Inches(1.9), y + Inches(0.08), Inches(2.5), Inches(0.65),
                name, font_size=20, bold=True, color=color)
    add_textbox(slide4, Inches(4.5), y + Inches(0.08), Inches(7.0), Inches(0.65),
                tech, font_size=15, color=LIGHT_GRAY)

add_textbox(slide4, Inches(1.5), Inches(6.8), Inches(10.3), Inches(0.5),
            '后端: Java 21  |  Spring Boot 3.2  |  MyBatis-Plus  |  JWT  |  BCrypt  |  AOP 审计日志',
            font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

print("Slide 4 done", flush=True)

# --- Slide 5: Features ---
add_section_slide(prs, 3, '六大功能模块')
slide5 = prs.slides[-1]

modules = [
    ('F01', '用户管理', '注册/登录/个人中心\n管理员用户CRUD', ACCENT_BLUE),
    ('F02', '文生图', '文本→图像AI生成\n参数调整/风格选择/历史', ACCENT_PURPLE),
    ('F03', '图生文', '图像→文本智能分析\n结果编辑/历史查看', ACCENT_GREEN),
    ('F04', '风格管理', '风格创建/编辑/删除\n预览展示/效果演示', ACCENT_ORANGE),
    ('F05', 'LoRA训练', '数据集上传/参数配置\n训练监控/模型下载', ACCENT_RED),
    ('F06', '系统管理', '用户管理/系统配置\n审计日志/模型管理', RGBColor(0x7C, 0xE2, 0xFE)),
]
for i, (code, title, desc, color) in enumerate(modules):
    col = i % 3
    row = i // 3
    left = Inches(1.5 + col * 3.6)
    top = Inches(1.5 + row * 2.9)
    card = add_rounded_rect(slide5, left, top, Inches(3.2), Inches(2.6), DARK_CARD)
    acc = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.2), Inches(0.06))
    acc.fill.solid()
    acc.fill.fore_color.rgb = color
    acc.line.fill.background()
    add_textbox(slide5, left + Inches(0.25), top + Inches(0.3), Inches(2.7), Inches(0.5),
                code, font_size=14, color=color, font_name='Consolas')
    add_textbox(slide5, left + Inches(0.25), top + Inches(0.7), Inches(2.7), Inches(0.5),
                title, font_size=20, bold=True, color=WHITE)
    add_textbox(slide5, left + Inches(0.25), top + Inches(1.3), Inches(2.7), Inches(1.2),
                desc, font_size=14, color=LIGHT_GRAY)

print("Slide 5 done", flush=True)

# --- Slide 6: AI Core ---
add_section_slide(prs, 4, 'AI 策略模式 & ComfyUI 集成')
slide6 = prs.slides[-1]

# Left: Strategy Pattern
add_textbox(slide6, Inches(1.5), Inches(1.2), Inches(5.5), Inches(0.6),
            '🧩 AI Provider 策略模式', font_size=24, bold=True, color=ACCENT_BLUE)
add_textbox(slide6, Inches(1.5), Inches(1.9), Inches(5.8), Inches(3.2),
            'Text2ImageProvider (文生图接口)\n'
            '  ├─ DoubaoText2ImageProvider\n'
            '  ├─ ComfyUIText2ImageProvider\n'
            '  └─ SD1.5Text2ImageProvider\n\n'
            'Image2TextProvider (图生文接口)\n'
            '  └─ DoubaoImage2TextProvider\n\n'
            '▸ application.yml 配置切换\n'
            '▸ DynamicConfigService 运行时热更新\n'
            '▸ 前端 AiStatusBar 实时切换',
            font_size=15, color=LIGHT_GRAY, font_name='Consolas')

# Right: ComfyUI
add_textbox(slide6, Inches(8.0), Inches(1.2), Inches(5.0), Inches(0.6),
            '🖥️ ComfyUI 深度集成', font_size=24, bold=True, color=ACCENT_GREEN)
add_textbox(slide6, Inches(8.0), Inches(1.9), Inches(5.5), Inches(3.2),
            'ComfyUIClient (HTTP协议封装)\n'
            '  POST /prompt → 提交工作流\n'
            '  轮询 GET /history/{id}\n'
            '  GET /view → 获取生成图片\n\n'
            'Workflow映射:\n'
            '  节点26 → 正向提示词\n'
            '  节点22 → 负向提示词\n'
            '  节点24 → 采样参数\n'
            '  节点25 → 图片尺寸\n\n'
            '▸ 默认API: http://127.0.0.1:8188\n'
            '▸ 管理后台可动态配置地址',
            font_size=15, color=LIGHT_GRAY, font_name='Consolas')

# Bottom stats
add_rounded_rect(slide6, Inches(1.5), Inches(5.4), Inches(10.3), Inches(1.6), DARK_CARD)
add_textbox(slide6, Inches(1.8), Inches(5.6), Inches(9.8), Inches(1.2),
            '📊  单元测试: ComfyUIClientTest (13个) + ComfyUIText2ImageProviderTest (7个) = 共20个\n'
            '📁  默认工作流文件: backend/src/main/resources/workflow/comfyui_default.json\n'
            '📐  文件上传限制: 单文件50MB  |  总请求100MB  |  JWT过期时间24h\n'
            '⚡  AI Provider 支持运行时切换，无需重启服务',
            font_size=15, color=LIGHT_GRAY)

print("Slide 6 done", flush=True)

# --- Slide 7: RBAC ---
add_section_slide(prs, 5, 'RBAC 三级角色权限')
slide7 = prs.slides[-1]

roles = [
    ('普通用户 USER', '文生图\n图生文\n历史记录\n个人设置', '🔵', ACCENT_BLUE),
    ('设计师 DESIGNER', '普通用户权限 +\n风格管理\nLoRA训练\n模型管理', '🟣', ACCENT_PURPLE),
    ('管理员 ADMIN', '全部权限 +\n用户管理\n系统配置\n审计日志', '🟢', ACCENT_GREEN),
]
for i, (role, perms, icon, color) in enumerate(roles):
    left = Inches(1.5 + i * 3.7)
    add_rounded_rect(slide7, left, Inches(1.5), Inches(3.3), Inches(3.8), DARK_CARD)
    c = slide7.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.15), Inches(1.8), Inches(1.0), Inches(1.0))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    # icon text on circle
    add_textbox(slide7, left + Inches(1.15), Inches(1.95), Inches(1.0), Inches(0.8),
                icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_textbox(slide7, left + Inches(0.2), Inches(3.0), Inches(2.9), Inches(0.6),
                role, font_size=20, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide7, left + Inches(0.3), Inches(3.6), Inches(2.7), Inches(1.5),
                perms, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Security features
add_rounded_rect(slide7, Inches(1.5), Inches(5.8), Inches(10.3), Inches(1.2), DARK_CARD)
add_textbox(slide7, Inches(1.8), Inches(6.0), Inches(9.8), Inches(0.8),
            '🔒 JWT Token 认证 (24h过期)  |  BCrypt 密码哈希  |  前端 Vue Router 路由守卫\n'
            '📝 AOP 审计日志 (audit_log表)  |  RBAC 路由级控制  |  默认管理员: admin@nebula.com / admin123',
            font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

print("Slide 7 done", flush=True)

# --- Slide 8: Database ---
add_section_slide(prs, 6, '数据库设计')
slide8 = prs.slides[-1]

tables = [
    ('user', '用户表', 'id, username, email, password, role, avatar, status, created_at', ACCENT_BLUE),
    ('image_record', '图像记录表', 'id, user_id, prompt, image_url, type, style_id, params, created_at', ACCENT_PURPLE),
    ('style', '风格表', 'id, name, description, preview_url, config, creator_id, created_at', ACCENT_GREEN),
    ('training_task', '训练任务表', 'id, name, status, params, dataset_path, output_path, creator_id', ACCENT_ORANGE),
    ('lora_model', 'LoRA模型表', 'id, name, version, file_path, training_task_id, style_id', ACCENT_RED),
    ('system_config', '系统配置表', 'id, config_key, config_value, description, updated_at', RGBColor(0x7C, 0xE2, 0xFE)),
    ('audit_log', '审计日志表', 'id, user_id, action, target, detail, ip, created_at', LIGHT_GRAY),
]
for i, (tbl, name, columns, color) in enumerate(tables):
    y = Inches(1.3 + i * 0.8)
    add_rounded_rect(slide8, Inches(1.5), y, Inches(10.3), Inches(0.68), DARK_CARD)
    add_textbox(slide8, Inches(1.7), y + Inches(0.08), Inches(0.5), Inches(0.5),
                f'0{i+1}', font_size=16, bold=True, color=color, font_name='Arial')
    add_textbox(slide8, Inches(2.3), y + Inches(0.08), Inches(2.0), Inches(0.5),
                tbl, font_size=16, bold=True, color=WHITE, font_name='Consolas')
    add_textbox(slide8, Inches(4.3), y + Inches(0.08), Inches(2.2), Inches(0.5),
                name, font_size=14, color=color)
    add_textbox(slide8, Inches(6.5), y + Inches(0.08), Inches(5.3), Inches(0.5),
                columns, font_size=12, color=LIGHT_GRAY, font_name='Consolas')

add_textbox(slide8, Inches(1.5), Inches(7.0), Inches(10.3), Inches(0.4),
            'H2 (开发/测试)  ⇄  MySQL 8.0+ (生产)  |  DataInitializer 启动时自动建表并初始化  |  MyBatis-Plus ORM',
            font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

print("Slide 8 done", flush=True)

# --- Slide 9: Frontend ---
add_section_slide(prs, 7, '前端架构一览')
slide9 = prs.slides[-1]

add_textbox(slide9, Inches(1.5), Inches(1.2), Inches(5.5), Inches(0.6),
            '🧱 Vue 3 组件结构', font_size=22, bold=True, color=ACCENT_BLUE)
add_textbox(slide9, Inches(1.5), Inches(1.9), Inches(5.8), Inches(4.8),
            'App.vue\n'
            '├── AppHeader / AppSidebar / AiStatusBar\n'
            '├── LoginView / RegisterView\n'
            '├── DashboardView\n'
            '├── Text2ImageView (文生图)\n'
            '├── Image2TextView (图生文)\n'
            '├── StyleManagementView\n'
            '├── LoraTrainingView\n'
            '├── AdminDashboard\n'
            '├── SystemConfigView\n'
            '├── UserManagementView\n'
            '└── AuditLogView',
            font_size=14, color=LIGHT_GRAY, font_name='Consolas')

# Right
add_textbox(slide9, Inches(8.0), Inches(1.2), Inches(5), Inches(0.6),
            '🛠️ 前端技术栈', font_size=22, bold=True, color=ACCENT_GREEN)
techs = [
    ('Vue 3 + Composition API', '核心框架'),
    ('Vite', '极速构建工具'),
    ('Element Plus', '企业级UI组件库'),
    ('Pinia', '轻量状态管理'),
    ('Vue Router', '路由守卫 + RBAC'),
    ('Axios', 'HTTP客户端 (request.js封装)'),
    ('VueUse', '组合式工具函数'),
]
for i, (tech, desc) in enumerate(techs):
    y = Inches(1.9 + i * 0.6)
    add_textbox(slide9, Inches(8.0), y, Inches(2.8), Inches(0.5),
                tech, font_size=16, bold=True, color=ACCENT_GREEN)
    add_textbox(slide9, Inches(10.8), y, Inches(1.8), Inches(0.5),
                desc, font_size=14, color=LIGHT_GRAY)

add_rounded_rect(slide9, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.7), DARK_CARD)
add_textbox(slide9, Inches(1.8), Inches(6.55), Inches(9.8), Inches(0.5),
            '端口 3000  |  API baseURL: http://localhost:8080/api  |  request.js 统一封装  |  路由级权限控制',
            font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

print("Slide 9 done", flush=True)

# --- Slide 10: Development ---
add_section_slide(prs, 8, '开发管理与部署')
slide10 = prs.slides[-1]

add_textbox(slide10, Inches(1.5), Inches(1.2), Inches(5.5), Inches(0.6),
            '📋 架构决策记录 (ADR)', font_size=22, bold=True, color=ACCENT_BLUE)
adrs = [
    ('ADR-001', 'Java 21 + Spring Boot 3.2 技术栈'),
    ('ADR-002', 'MySQL 8.0+ 数据库 (生产环境)'),
    ('ADR-003', 'Redis 缓存方案'),
    ('ADR-004', 'Kafka 消息队列'),
    ('ADR-005', 'MinIO 文件存储'),
    ('ADR-006', 'Vue 3 前端框架'),
    ('ADR-007', 'JWT 认证方案'),
    ('ADR-008', 'RBAC 授权方案'),
    ('ADR-009', 'AI Provider 策略模式 + ComfyUI 集成'),
]
for i, (code, desc) in enumerate(adrs):
    y = Inches(1.9 + i * 0.45)
    add_textbox(slide10, Inches(1.8), y, Inches(1.0), Inches(0.4),
                code, font_size=13, bold=True, color=ACCENT_BLUE, font_name='Consolas')
    add_textbox(slide10, Inches(2.8), y, Inches(4.5), Inches(0.4),
                '✅ ' + desc, font_size=13, color=LIGHT_GRAY)

# Right: Deployment
add_textbox(slide10, Inches(8.0), Inches(1.2), Inches(5), Inches(0.6),
            '🚀 部署命令', font_size=22, bold=True, color=ACCENT_GREEN)
add_textbox(slide10, Inches(8.0), Inches(1.9), Inches(5.5), Inches(4.2),
            '开发环境 (H2):\n'
            '  build.bat     → 构建项目\n'
            '  start.bat     → 启动服务\n'
            '  stop.bat      → 停止服务\n'
            '  dev.bat       → 管理面板\n'
            '  test.bat      → 冒烟测试\n\n'
            '生产环境 (MySQL):\n'
            '  mvn spring-boot:run\n'
            '    --spring.profiles.active=mysql\n\n'
            '容器化:\n'
            '  Docker Compose + K8s 编排',
            font_size=14, color=LIGHT_GRAY, font_name='Consolas')

# Bottom
add_rounded_rect(slide10, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.7), DARK_CARD)
add_textbox(slide10, Inches(1.8), Inches(6.55), Inches(9.8), Inches(0.5),
            '📦 20个单元测试  |  Python全量API测试  |  前后端独立部署  |  H2开发 / MySQL生产',
            font_size=15, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

print("Slide 10 done", flush=True)

# --- Slide 11: End ---
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide11)
add_textbox(slide11, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.5),
            '感谢聆听', font_size=56, bold=True, font_name='Arial')
add_textbox(slide11, Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.0),
            'Q & A', font_size=32, color=ACCENT_BLUE, font_name='Arial')
line = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.8), Inches(3), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_PURPLE
line.line.fill.background()
add_textbox(slide11, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.6),
            'Nebula Studio — 文图互转主题设计系统', font_size=20, color=LIGHT_GRAY)
add_textbox(slide11, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.6),
            '默认管理员: admin@nebula.com  |  前端端口: 3000  |  后端端口: 8080', font_size=16, color=LIGHT_GRAY)

print("Slide 11 done", flush=True)

# ==================== SAVE ====================
output_path = r'C:\Users\syq04\Downloads\project01\Nebula_Studio_项目汇报.pptx'
prs.save(output_path)
print(f'PPT saved to: {output_path}', flush=True)
print('Done!')
